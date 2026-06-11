from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_slide(prs, layout_idx, title, bullets):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
        
    if len(slide.shapes) > 1 and slide.shapes[1].has_text_frame:
        tf = slide.shapes[1].text_frame
        tf.clear()
        
        for i, b in enumerate(bullets):
            p = tf.add_paragraph()
            p.text = b
            if i == 0:
                p.level = 0
            else:
                if b.startswith("    "):
                    p.level = 2
                elif b.startswith("  "):
                    p.level = 1
                else:
                    p.level = 0
            
            p.font.size = Pt(20)

def main():
    prs = Presentation()

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Driven Environmental Monitoring and Shelf-Life Prediction System"
    subtitle.text = "Final Semester Project Defense\nFood Spoilage Detector\n\n[Add Team Members / IDs]"

    # Slide 2
    add_slide(prs, 1, "Mid-Sem Recap & Final Objectives", [
        "The Problem: Post-harvest loss due to poor storage conditions.",
        "Our Objective: Move beyond simple sensor readings to predictive AI.",
        "What We Built (End-Sem Focus):",
        "  • Physics-grounded synthetic data generator.",
        "  • Dual-model Machine Learning Engine (Random Forest + LSTM).",
        "  • Two-stage environmental rules engine.",
        "  • Premium SaaS-style monitoring dashboard."
    ])

    # Slide 3
    add_slide(prs, 1, "Overall System Architecture", [
        "Data Ingestion:",
        "  • IoT Sensors capture Temp, RH, CO2, Methane, and Ethylene.",
        "Backend:",
        "  • High-performance Python FastAPI server handling data streams and CSV batch uploads.",
        "Intelligence:",
        "  • Scikit-Learn (Random Forest) and TensorFlow (LSTM) engines predicting Remaining Shelf Life (RSL).",
        "User Interface:",
        "  • React + Tailwind dashboard for real-time visualization."
    ])

    # Slide 4
    add_slide(prs, 1, "Data Generation — The Physics Simulator", [
        "The Challenge:",
        "  • Collecting 100,000+ real-world sensor rows is impossible in one semester.",
        "The Solution:",
        "  • We built a strict, physics-grounded simulator to generate training data.",
        "Core Mathematics:",
        "  • Arrhenius Factor: Models exponential biological aging based on temperature.",
        "  • Climacteric Multiplier: Models the sudden spike in respiration and Ethylene during ripening.",
        "Result:",
        "  • Generated over 105,000 rows of highly realistic, scientifically accurate training data."
    ])

    # Slide 5
    add_slide(prs, 1, "Machine Learning Architecture", [
        "We evaluated two distinct models to predict Remaining Shelf Life (RSL):",
        "Model 1: Random Forest (Baseline)",
        "  • Treats every sensor reading independently.",
        "  • Fast and highly interpretable.",
        "Model 2: LSTM (Long Short-Term Memory Network)",
        "  • A deep learning model trained on a 48-step sliding window.",
        "Why LSTM?",
        "  • Food spoilage is a time-series problem. How the gas accumulated over the last 24 hours is just as important as the current reading."
    ])

    # Slide 6
    add_slide(prs, 1, "Model Performance & Results", [
        "Random Forest Performance:",
        "  • RMSE = 0.077 days.",
        "LSTM Performance:",
        "  • RMSE = 0.043 days (R-squared = 0.9997).",
        "Key Finding:",
        "  • The LSTM almost halved the error rate.",
        "Conclusion:",
        "  • By utilizing the time-history representation (memory of past sensor readings), the LSTM suppresses large errors and provides a highly stable, production-ready prediction."
    ])

    # Slide 7
    add_slide(prs, 1, "Environmental Alert Engine (Traffic Light Logic)", [
        "Two-Stage Decision Logic:",
        "  • The system actively monitors the environment for early warnings, not just total spoilage.",
        "Stage 1: Sensor Thresholds:",
        "  • Critical: Methane > 0.5 ppm (Anaerobic rot started).",
        "  • Warning: Temp > 30°C or CO2 > 50,000 ppm (Dangerously hot/suffocating).",
        "Stage 2: Aggregation:",
        "  • Factors are combined with the ML prediction to assign a final actionable status (Green/Raw, Yellow/Ripe, Red/Rotten)."
    ])

    # Slide 8
    add_slide(prs, 1, "The Operational Dashboard", [
        "We designed a premium, startup-level user interface.",
        "Tech Stack:",
        "  • React, TypeScript, Tailwind CSS (Glassmorphism design), Recharts, and Framer Motion.",
        "Features:",
        "  • Real-time glowing sensor cards (Temp, RH, Gases).",
        "  • Dynamic Area Charts mapping the distribution trends.",
        "  • A giant countdown timer for Remaining Shelf Life (RSL).",
        "[Insert Screenshot of Live Dashboard Here]"
    ])

    # Slide 9
    add_slide(prs, 1, "Batch Processing & Analytics", [
        "Enterprise Capability:",
        "  • Built a batch-processing engine for warehouse managers.",
        "Functionality:",
        "  • Users can drag-and-drop massive CSV files of raw sensor data.",
        "Analysis:",
        "  • The LSTM engine processes thousands of rows concurrently.",
        "  • Provides a breakdown of dataset status percentages (Raw vs Rotten).",
        "  • Per-row inspections identifying exact confidence and reasons for spoilage.",
        "[Insert Screenshot of CSV Upload Panel Here]"
    ])

    # Slide 10
    add_slide(prs, 1, "Conclusion & Future Scope", [
        "What We Achieved:",
        "  • Successfully mapped raw chemical gas readings to a highly accurate biological timeline using Deep Learning.",
        "Future Scope:",
        "  • Integrate physical ESP32/Raspberry Pi hardware directly into the WebSocket stream.",
        "  • Train the models on different climacteric fruits (e.g., Avocados, Mangoes).",
        "Questions?",
        "  • Thank you!"
    ])

    prs.save('Food_Spoilage_Detector_Presentation.pptx')
    print("Saved Presentation: Food_Spoilage_Detector_Presentation.pptx")

if __name__ == '__main__':
    main()
